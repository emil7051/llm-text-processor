import os
import re
import pandas as pd
import tiktoken
from pathlib import Path
import warnings

# Suppress urllib3 warnings
warnings.filterwarnings("ignore", category=UserWarning, module="urllib3")

def extract_text_from_file(file_path):
    """Extract text from various file formats based on extension."""
    ext = file_path.suffix.lower()
    
    try:
        # PDF Files
        if ext == '.pdf':
            import PyPDF2
            text = ""
            with open(file_path, "rb") as file:
                reader = PyPDF2.PdfReader(file)
                for page in reader.pages:
                    extracted = page.extract_text() or ""
                    text += extracted + "\n"
            return text, len(reader.pages)
            
        # Markdown/Text Files
        elif ext in ['.md', '.txt', '.rtf']:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as file:
                text = file.read()
            return text, 1  # Count as 1 page for simplicity
            
        # PowerPoint Files
        elif ext in ['.ppt', '.pptx']:
            from pptx import Presentation
            text = ""
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text, len(prs.slides)
            
        # Excel Files
        elif ext in ['.xls', '.xlsx', '.csv']:
            import pandas as pd
            
            if ext == '.csv':
                df = pd.read_csv(file_path, encoding='utf-8', errors='replace')
            else:
                df = pd.read_excel(file_path)
                
            # Convert all dataframe content to string and join
            text = "\n".join(df.astype(str).to_csv(index=False).splitlines())
            return text, len(df)  # Use row count as "page" count for spreadsheets
            
        # Word Documents
        elif ext in ['.doc', '.docx']:
            import docx
            doc = docx.Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
            return text, len(doc.paragraphs) // 3 or 1  # Approximate page count
            
        else:
            print(f"Unsupported file type: {ext} for file {file_path}")
            return "", 0
            
    except Exception as e:
        print(f"Error processing {file_path}: {e}")
        return "", 0

def count_tokens(text, encoding_name="cl100k_base"):
    """Count tokens in a text string using the specified encoding."""
    try:
        encoding = tiktoken.get_encoding(encoding_name)
        tokens = encoding.encode(text)
        return len(tokens)
    except Exception as e:
        print(f"Error counting tokens: {e}")
        return 0

def count_tokens_in_folder(folder_path, encoding_name="cl100k_base", recursive=True, 
                          file_types=None):
    """Count tokens in all supported files in a folder."""
    results = []
    
    # Set default file types if none provided
    if file_types is None:
        file_types = ['.pdf', '.md', '.txt', '.rtf', '.ppt', '.pptx', 
                     '.xls', '.xlsx', '.csv', '.doc', '.docx']
    
    # Convert file_types to lowercase for case-insensitive matching
    file_types = [ft.lower() if ft.startswith('.') else f'.{ft.lower()}' for ft in file_types]
    
    # Determine which files to process
    if recursive:
        all_files = list(Path(folder_path).glob("**/*.*"))
    else:
        all_files = list(Path(folder_path).glob("*.*"))
    
    # Filter files by extension
    files_to_process = [f for f in all_files if f.suffix.lower() in file_types]
    
    # Process each file
    total_tokens = 0
    total_files_processed = 0
    
    print(f"Found {len(files_to_process)} files to process...")
    
    for i, file_path in enumerate(files_to_process):
        print(f"Processing {i+1}/{len(files_to_process)}: {file_path.name}")
        
        text, page_count = extract_text_from_file(file_path)
        
        if not text:
            print(f"  Skipping {file_path.name} - no text extracted")
            continue
            
        token_count = count_tokens(text, encoding_name)
        
        # Calculate tokens per page
        tokens_per_page = token_count / page_count if page_count > 0 else 0
        
        # Determine file type for reporting
        file_type = file_path.suffix.lower().lstrip('.')
        
        # Store results
        results.append({
            "file": str(file_path),
            "type": file_type,
            "tokens": token_count,
            "pages/units": page_count,
            "tokens_per_unit": round(tokens_per_page, 1)
        })
        
        total_tokens += token_count
        total_files_processed += 1
    
    # Create summary DataFrame
    if results:
        df = pd.DataFrame(results)
        
        # Add type summary rows
        type_summary = df.groupby('type').agg({
            'tokens': 'sum',
            'pages/units': 'sum',
            'file': 'count'
        }).reset_index()
        
        type_summary = type_summary.rename(columns={'file': 'count'})
        type_summary['tokens_per_unit'] = round(type_summary['tokens'] / type_summary['pages/units'], 1)
        
        # Format the type summary rows
        for idx, row in type_summary.iterrows():
            type_summary.at[idx, 'file'] = f"TOTAL ({row['type']})"
        
        # Add grand total row
        grand_total = pd.DataFrame([{
            'file': 'GRAND TOTAL',
            'type': 'all',
            'tokens': total_tokens,
            'pages/units': df['pages/units'].sum(),
            'tokens_per_unit': round(total_tokens / df['pages/units'].sum() if df['pages/units'].sum() > 0 else 0, 1),
            'count': total_files_processed
        }])
        
        # Create final summary dataframe with individual files and summaries
        summary_df = pd.concat([df, type_summary[['file', 'type', 'tokens', 'pages/units', 'tokens_per_unit', 'count']], grand_total], ignore_index=True)
        
        return summary_df
    else:
        return pd.DataFrame(columns=["file", "type", "tokens", "pages/units", "tokens_per_unit"])

# Example usage
if __name__ == "__main__":
    import argparse
    
    parser = argparse.ArgumentParser(description="Count tokens in various file types within a folder")
    parser.add_argument("folder", help="Folder containing files to process")
    parser.add_argument("--model", default="cl100k_base", help="Tokeniser model to use (default: cl100k_base for GPT-4/3.5)")
    parser.add_argument("--no-recursive", action="store_true", help="Don't search subdirectories")
    parser.add_argument("--output", help="Save results to CSV file")
    parser.add_argument("--file-types", nargs="+", default=None, 
                        help="File extensions to process (default: all supported types)")
    
    args = parser.parse_args()
    
    print(f"Starting token count in folder: {args.folder}")
    print(f"Using tokeniser model: {args.model}")
    
    results = count_tokens_in_folder(
        args.folder, 
        encoding_name=args.model,
        recursive=not args.no_recursive,
        file_types=args.file_types
    )
    
    # Print results
    print(f"\nProcessed {len(results)-1} files")
    print(f"Total tokens: {results.loc[results['file'] == 'GRAND TOTAL', 'tokens'].values[0]:,}")
    print("\nSummary by file type:")
    type_summary = results[results['file'].str.contains('TOTAL')]
    print(type_summary[['file', 'tokens', 'count']].to_string(index=False))
    
    print("\nTop 10 files by token count:")
    individual_files = results[~results['file'].str.contains('TOTAL')]
    print(individual_files.sort_values('tokens', ascending=False).head(10)[['file', 'type', 'tokens']].to_string(index=False))
    
    # Save to CSV if requested
    if args.output:
        results.to_csv(args.output, index=False)
        print(f"\nResults saved to {args.output}")