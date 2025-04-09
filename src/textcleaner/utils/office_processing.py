"""Utilities for processing specific Office file formats."""

import re
from pathlib import Path
from typing import Any, Dict, Tuple

import docx
from docx.document import Document as DocxDocument
import pandas as pd
from pptx import Presentation

from textcleaner.utils.logging_config import get_logger

logger = get_logger(__name__)

def process_docx(file_path: Path, doc: DocxDocument) -> Tuple[str, Dict[str, Any]]:
    """Extracts text and metadata from a python-docx Document object.
    
    Args:
        file_path: Original path of the DOCX file (used for title fallback).
        doc: The python-docx Document object.
        
    Returns:
        Tuple of (extracted_text, metadata_dict).
    """
    text_parts = []
    metadata = {
        "paragraph_count": len(doc.paragraphs),
        "section_count": len(doc.sections),
    }

    # Try to extract core properties
    try:
        core_props = doc.core_properties
        if core_props.title:
            metadata["title"] = core_props.title
        if core_props.author:
            metadata["author"] = core_props.author
        if core_props.last_modified_by:
            metadata["last_modified_by"] = core_props.last_modified_by
    except Exception as e: # Catch potential attribute errors or other issues
        logger.warning(f"Could not extract core properties from {file_path}: {e}")
        
    # Add title if available
    if "title" in metadata:
        text_parts.append(f"# {metadata['title']}\n")
    else:
        # Use filename as title if no title in metadata
        text_parts.append(f"# {file_path.stem}\n")
    
    # Process paragraphs, headings, and lists
    in_list = False
    current_heading_level = 0
    
    for paragraph in doc.paragraphs:
        if not paragraph.text.strip():
            text_parts.append("") # Preserve paragraph breaks
            continue
            
        style_name = paragraph.style.name
        if style_name.startswith('Heading'):
            try:
                heading_level = int(style_name[-1])
                prefix = '#' * heading_level
                text_parts.append(f"\n{prefix} {paragraph.text}\n")
                current_heading_level = heading_level
                in_list = False
            except (ValueError, IndexError):
                 # Fallback for unexpected heading style names
                 text_parts.append(f"\n## {paragraph.text}\n") 
            continue
            
        if style_name.startswith('List'):
            if not in_list:
                text_parts.append("") # Add newline before list start
                in_list = True
            text_parts.append(f"* {paragraph.text}")
            continue
        
        # Regular paragraph
        in_list = False
        text_parts.append(paragraph.text)
        
    # Extract tables
    table_count = 0
    for table in doc.tables:
        table_count += 1
        text_parts.append(f"\n### Table {table_count}\n")
        
        rows = []
        if not table.rows:
            continue
            
        num_cols = len(table.columns)
        if num_cols == 0:
             continue
             
        # Assuming first row is header
        header_cells = table.rows[0].cells
        header = " | ".join(cell.text.strip() for cell in header_cells)
        rows.append(f"| {header} |")
        rows.append(f"| {' | '.join(['---'] * num_cols)} |")
        
        # Data rows
        for row in table.rows[1:]:
            data_cells = row.cells
            data = " | ".join(cell.text.strip() for cell in data_cells)
            rows.append(f"| {data} |")
        
        text_parts.append("\n".join(rows))
        text_parts.append("\n")
            
    metadata["table_count"] = table_count
            
    return "\n".join(text_parts), metadata

def process_excel(file_path: Path, 
                  excel_file: pd.ExcelFile, 
                  max_rows: int, 
                  max_cols: int) -> Tuple[str, Dict[str, Any]]:
    """Extracts text and metadata from a pandas ExcelFile object.
    
    Args:
        file_path: Original path of the Excel file (used for title).
        excel_file: The pandas ExcelFile object.
        max_rows: Maximum rows to extract per sheet.
        max_cols: Maximum columns to extract per sheet.
        
    Returns:
        Tuple of (extracted_text, metadata_dict).
    """
    text_parts = [f"# {file_path.name}\n"]
    sheet_count = len(excel_file.sheet_names)
    metadata = {
        "sheet_count": sheet_count,
        "sheets": {}
    }
    
    for sheet_index, sheet_name in enumerate(excel_file.sheet_names):
        clean_name = re.sub(r'[^\w\s.-]', '_', sheet_name)
        
        if sheet_count > 1:
            text_parts.append(f"## Sheet: {clean_name}\n")
        
        try:
            df = pd.read_excel(excel_file, sheet_name=sheet_name)
            sheet_metadata = {
                "row_count": len(df),
                "column_count": len(df.columns),
            }
            metadata["sheets"][clean_name] = sheet_metadata
            
            # Apply limits
            rows_truncated = False
            cols_truncated = False
            if len(df) > max_rows:
                df = df.head(max_rows)
                rows_truncated = True
                
            if len(df.columns) > max_cols:
                df = df.iloc[:, :max_cols]
                cols_truncated = True
                
            if rows_truncated:
                 text_parts.append(f"*Note: Table truncated to {max_rows} rows*\n")
            if cols_truncated:
                 text_parts.append(f"*Note: Table truncated to {max_cols} columns*\n")
            
            if df.empty:
                text_parts.append("*Empty sheet*\n")
                continue
                
            df = df.fillna("")
            df.columns = df.columns.astype(str)
            
            header = "| " + " | ".join(str(col) for col in df.columns) + " |"
            separator = "| " + " | ".join(["---"] * len(df.columns)) + " |"
            
            rows = []
            for _, row in df.iterrows():
                row_str = "| " + " | ".join(str(val).replace("\n", "<br>") for val in row.values) + " |"
                rows.append(row_str)
            
            text_parts.append(header)
            text_parts.append(separator)
            text_parts.extend(rows)
            text_parts.append("\n")
            
        except pd.errors.ParserError as sheet_error:
            logger.warning(f"Pandas parsing error on sheet '{clean_name}' in {file_path}: {sheet_error}")
            text_parts.append(f"*Error parsing sheet {clean_name}*\n")
        except Exception as sheet_error:
            logger.exception(f"Unexpected error reading sheet '{clean_name}' in {file_path}")
            text_parts.append(f"*Unexpected error reading sheet {clean_name}: {sheet_error}*\n")
        
        if sheet_index < sheet_count - 1:
            text_parts.append("---\n")
            
    return "\n".join(text_parts), metadata 

def process_pptx(file_path: Path, prs: Presentation) -> Tuple[str, Dict[str, Any]]:
    """Extracts text and metadata from a python-pptx Presentation object.
    
    Args:
        file_path: Original path of the PPTX file (used for title).
        prs: The python-pptx Presentation object.
        
    Returns:
        Tuple of (extracted_text, metadata_dict).
    """
    text_parts = [f"# {file_path.name}\n"]
    metadata = {
        "slide_count": len(prs.slides),
    }
    
    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        text_parts.append(f"\n## Slide {slide_num}")
        
        slide_title = None
        if slide.shapes.title and slide.shapes.title.has_text_frame and slide.shapes.title.text:
            slide_title = slide.shapes.title.text.strip()
            text_parts.append(f": {slide_title}\n")
        else:
            text_parts.append("\n")
        
        shape_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text:
                # Skip the title shape if we processed it already
                if shape == slide.shapes.title and slide_title:
                    continue
                shape_texts.append(shape.text.strip())
        
        if len(shape_texts) > 1:
            for text in shape_texts:
                # Simple formatting: treat each shape text as a bullet point
                cleaned_text = text.replace('\n', '<br>')
                text_parts.append(f"* {cleaned_text}")
        elif shape_texts:
             cleaned_text = shape_texts[0].replace('\n', '<br>')
             text_parts.append(cleaned_text)
        
        # Add notes if present
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                 text_parts.append(f"\n**Notes:**\n{notes}")
                 metadata.setdefault("slides_with_notes", []).append(slide_num)
                 
        # Separator
        if i < len(prs.slides) - 1:
             text_parts.append("\n---\n")
        
    return "\n".join(text_parts), metadata 